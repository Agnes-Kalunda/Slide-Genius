import os
import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from django.conf import settings
from django.shortcuts import render
from django.http import JsonResponse
from PIL import Image, ImageDraw, ImageFont

openai.api_key = os.getenv('OPENAI_API_KEY')

def generate_presentation(topic, num_slides):
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": f"Create a summary presentation about '{topic}' consisting of {num_slides} slides with titles and bullet points for just the content."}
            ],
            max_tokens=600,
            temperature=0.1
        )
        slides_content = response.choices[0].message['content'].strip().split('\n\n')

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]

        # Ensure the media directory exists
        media_dir = settings.MEDIA_ROOT
        if not os.path.exists(media_dir):
            os.makedirs(media_dir)

        # Load a font
        font_path = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
        title_font = ImageFont.truetype(font_path, 40)  # Increase font size for title
        content_font = ImageFont.truetype(font_path, 20)  # Increase font size for content

        slide_image_paths = []

        title_styles = [
            {'title_color': (255, 0, 0)},  # Red
            {'title_color': (0, 0, 255)},  # Blue
            {'title_color': (0, 128, 0)},  # Green
            {'title_color': (255, 165, 0)},  # Orange
        ]

        for i, slide_content in enumerate(slides_content):
            if ':' in slide_content:
                title, content = slide_content.split(':', 1)
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = f"Slide {i + 1}"

                content_placeholder = slide.placeholders[1]
                content_frame = content_placeholder.text_frame
                content_frame.word_wrap = True  # Enable word wrap

                p = content_frame.add_paragraph()
                p.text = title.strip()
                p.font.bold = True  # Bold the title
                p.font.size = Pt(24)  # Adjust title font size
                p.level = 0

                for line in content.strip().split('\n'):
                    p = content_frame.add_paragraph()
                    p.text = line.strip()
                    p.font.size = Pt(18)  # Adjust content font size
                    p.level = 1  # Indent bullet points

                #  text frame size
                left = Inches(1)
                top = Inches(1.5)
                width = Inches(8.5)
                height = Inches(5)
                content_placeholder.width = width
                content_placeholder.height = height
                content_placeholder.left = left
                content_placeholder.top = top

                # title styles
                title_style = title_styles[i % len(title_styles)]
                image_path = os.path.join(media_dir, f"{topic}_slide_{i + 1}.png")
                img = Image.new('RGB', (800, 600), color=(255, 255, 255))  # White background
                d = ImageDraw.Draw(img)
                d.text((10, 10), title.strip(), font=title_font, fill=title_style['title_color'])

                y_position = 100
                for line in content.strip().split('\n'):
                    d.text((10, y_position), f"- {line.strip()}", font=content_font, fill=(0, 0, 0))
                    y_position += 40  # Adjust line height

                img.save(image_path)
                slide_image_paths.append(image_path.replace(settings.MEDIA_ROOT, '').lstrip('/'))

        prs_path = os.path.join(media_dir, f"{topic}.pptx")
        prs.save(prs_path)

        return slide_image_paths
    except Exception as e:
        raise e

def chat_view(request):
    if request.method == 'POST':
        user_input = request.POST.get('message', '').strip()

        if not user_input:
            return JsonResponse({'error': 'Please enter a message'})

        if 'topic' not in request.session or request.session['topic'] is None:
            request.session['topic'] = user_input
            return JsonResponse({'chatgpt_response': f"How many slides would you like for the topic: '{user_input}'?"})
        else:
            topic = request.session['topic']
            if user_input.isdigit() and 1 <= int(user_input) <= 20:
                num_slides = int(user_input)
                try:
                    slide_image_paths = generate_presentation(topic, num_slides)
                    request.session.pop('topic')
                    return JsonResponse({'chatgpt_response': f"Here is your presentation on '{topic}' . Provide a new topic if you want a different presentation. ", 'slide_image_paths': slide_image_paths})
                except Exception as e:
                    return JsonResponse({'error': f'Failed to generate presentation: {str(e)}'})
            else:
                return JsonResponse({'error': 'Please enter a number of slides between 1-20'})

    return render(request, 'index.html')
